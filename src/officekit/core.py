"""officekit - friendly Office document automation for humans and AI agents.

A small, well-typed toolkit to *read*, *create*, and *fill* Office documents
(.docx / .xlsx / .pptx) either from the command line or as a library.

This is an original implementation in the spirit of agent-facing Office CLIs:
it favours a predictable, scriptable surface over a sprawling feature set.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Dict, List

from docx import Document as DocxDocument
from openpyxl import Workbook, load_workbook
from pptx import Presentation

# {{ key }}  or  {{key}}  — whitespace tolerant
PLACEHOLDER = re.compile(r"\{\{\s*([\w.\-]+)\s*\}\}")

_KIND_EXT = {"docx": ".docx", "xlsx": ".xlsx", "pptx": ".pptx"}


def _detect_kind(path: str | Path) -> str:
    suffix = Path(path).suffix.lower()
    if suffix == ".docx":
        return "docx"
    if suffix == ".xlsx":
        return "xlsx"
    if suffix == ".pptx":
        return "pptx"
    raise ValueError(
        f"Unsupported file type: {suffix!r} (expected .docx / .xlsx / .pptx)"
    )


def _fill(text: str, data: Dict[str, Any]) -> str:
    """Replace {{ key }} tokens with values from *data* (unknown keys kept)."""
    return PLACEHOLDER.sub(lambda m: str(data.get(m.group(1), m.group(0))), text)


# --------------------------------------------------------------------------
# extract_text
# --------------------------------------------------------------------------
def extract_text(path: str | Path) -> str:
    """Return the visible text of a document as a single string."""
    kind = _detect_kind(path)
    if kind == "docx":
        return _extract_docx(path)
    if kind == "xlsx":
        return _extract_xlsx(path)
    return _extract_pptx(path)


def _extract_docx(path: str | Path) -> str:
    doc = DocxDocument(str(path))
    parts: List[str] = []
    for p in doc.paragraphs:
        if p.text:
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text for c in row.cells]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_xlsx(path: str | Path) -> str:
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(vals):
                parts.append("\t".join(vals))
    return "\n".join(parts)


def _extract_pptx(path: str | Path) -> str:
    prs = Presentation(str(path))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text:
                        parts.append(p.text)
    return "\n".join(parts)


# --------------------------------------------------------------------------
# create_document
# --------------------------------------------------------------------------
def create_document(
    kind: str, path: str | Path, title: str | None = None
) -> Path:
    """Create a starter document of *kind* (docx/xlsx/pptx) at *path*."""
    kind = kind.lower()
    if kind not in _KIND_EXT:
        raise ValueError(f"Unknown kind: {kind!r} (expected docx/xlsx/pptx)")
    path = Path(path)
    if path.suffix.lower() != _KIND_EXT[kind]:
        path = path.with_suffix(_KIND_EXT[kind])
    if kind == "docx":
        _create_docx(path, title)
    elif kind == "xlsx":
        _create_xlsx(path, title)
    else:
        _create_pptx(path, title)
    return path


def _create_docx(path: Path, title: str | None) -> None:
    doc = DocxDocument()
    doc.add_heading(title or "Untitled Document", level=0)
    doc.add_paragraph("Created with officekit.")
    doc.save(str(path))


def _create_xlsx(path: Path, title: str | None) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = title or "Sheet1"
    ws["A1"] = title or "officekit"
    wb.save(str(path))


def _create_pptx(path: Path, title: str | None) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = title or "officekit"
    prs.save(str(path))


# --------------------------------------------------------------------------
# merge_template
# --------------------------------------------------------------------------
def merge_template(
    path: str | Path, out: str | Path, data: Dict[str, Any]
) -> Path:
    """Fill {{ key }} placeholders in *path* and write the result to *out*."""
    kind = _detect_kind(path)
    out = Path(out)
    if out.suffix.lower() != _KIND_EXT[kind]:
        out = out.with_suffix(_KIND_EXT[kind])
    if kind == "docx":
        _merge_docx(path, out, data)
    elif kind == "xlsx":
        _merge_xlsx(path, out, data)
    else:
        _merge_pptx(path, out, data)
    return out


def _merge_docx(path: str | Path, out: Path, data: Dict[str, Any]) -> None:
    doc = DocxDocument(str(path))
    for p in doc.paragraphs:
        for run in p.runs:
            if PLACEHOLDER.search(run.text):
                run.text = _fill(run.text, data)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for run in p.runs:
                        if PLACEHOLDER.search(run.text):
                            run.text = _fill(run.text, data)
    doc.save(str(out))


def _merge_xlsx(path: str | Path, out: Path, data: Dict[str, Any]) -> None:
    wb = load_workbook(str(path))
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and PLACEHOLDER.search(cell.value):
                    cell.value = _fill(cell.value, data)
    wb.save(str(out))


def _merge_pptx(path: str | Path, out: Path, data: Dict[str, Any]) -> None:
    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if PLACEHOLDER.search(run.text):
                            run.text = _fill(run.text, data)
    prs.save(str(out))


# --------------------------------------------------------------------------
# document_info / validate_document
# --------------------------------------------------------------------------
def document_info(path: str | Path) -> Dict[str, Any]:
    """Return a summary dict of document shape (counts, sheet names, ...)."""
    kind = _detect_kind(path)
    if kind == "docx":
        doc = DocxDocument(str(path))
        return {
            "kind": "docx",
            "paragraphs": len([p for p in doc.paragraphs if p.text]),
            "tables": len(doc.tables),
        }
    if kind == "xlsx":
        wb = load_workbook(str(path), read_only=True)
        return {
            "kind": "xlsx",
            "sheets": len(wb.sheetnames),
            "sheet_names": wb.sheetnames,
        }
    prs = Presentation(str(path))
    return {"kind": "pptx", "slides": len(prs.slides)}


def validate_document(path: str | Path) -> List[str]:
    """Open the document and return a list of issues (empty == healthy)."""
    issues: List[str] = []
    path = Path(path)
    if not path.exists():
        return ["File does not exist"]
    try:
        kind = _detect_kind(path)
        if kind == "docx":
            DocxDocument(str(path))
        elif kind == "xlsx":
            load_workbook(str(path))
        else:
            Presentation(str(path))
    except Exception as exc:  # noqa: BLE001 - we surface the reason as text
        issues.append(f"Cannot open document: {exc}")
    return issues
