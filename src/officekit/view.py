"""Semantic views of a document — outline / stats / issues.

These mirror the ``view`` family of commands from agent-facing Office CLIs:
a quick, machine-readable way to *understand* a document before editing it.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from .core import PLACEHOLDER, _detect_kind


def view_outline(path: str | Path) -> Dict[str, Any]:
    """Return the document skeleton (headings, sheets, slide titles)."""
    kind = _detect_kind(path)
    if kind == "docx":
        doc = DocxDocument(str(path))
        headings: List[Dict[str, Any]] = []
        for p in doc.paragraphs:
            style = (p.style.name or "") if p.style else ""
            if style.startswith("Heading") and p.text.strip():
                level = 1
                digits = style.replace("Heading", "").strip()
                if digits.isdigit():
                    level = int(digits)
                headings.append({"level": level, "text": p.text})
        return {"kind": "docx", "headings": headings}
    if kind == "xlsx":
        wb = load_workbook(str(path), read_only=True)
        sheets = [
            {"name": name, "dimension": wb[name].dimensions}
            for name in wb.sheetnames
        ]
        return {"kind": "xlsx", "sheets": sheets}
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text or ""
        slides.append({"index": i, "title": title})
    return {"kind": "pptx", "slides": slides}


def view_stats(path: str | Path) -> Dict[str, Any]:
    """Return counts and dimensions for the document."""
    kind = _detect_kind(path)
    if kind == "docx":
        doc = DocxDocument(str(path))
        words = sum(len(p.text.split()) for p in doc.paragraphs if p.text)
        return {
            "kind": "docx",
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
            "words": words,
            "images": len(doc.inline_shapes),
        }
    if kind == "xlsx":
        wb = load_workbook(str(path), read_only=True)
        sheet_detail = [
            {"name": name, "rows": wb[name].max_row, "cols": wb[name].max_column}
            for name in wb.sheetnames
        ]
        return {
            "kind": "xlsx",
            "sheets": len(wb.sheetnames),
            "sheet_detail": sheet_detail,
        }
    prs = Presentation(str(path))
    return {
        "kind": "pptx",
        "slides": len(prs.slides),
        "shapes": sum(len(s.shapes) for s in prs.slides),
    }


def view_issues(path: str | Path) -> List[Dict[str, Any]]:
    """Return a list of potential problems (empty == healthy)."""
    kind = _detect_kind(path)
    issues: List[Dict[str, Any]] = []
    if kind == "docx":
        doc = DocxDocument(str(path))
        text = "\n".join(p.text for p in doc.paragraphs)
        if not text.strip() and not doc.tables:
            issues.append({"code": "empty", "message": "Document has no text content"})
        for m in PLACEHOLDER.finditer(text):
            issues.append(
                {
                    "code": "unfilled_placeholder",
                    "message": f"Unfilled placeholder {m.group(0)}",
                }
            )
    elif kind == "xlsx":
        wb = load_workbook(str(path), read_only=True)
        if not wb.sheetnames:
            issues.append({"code": "empty", "message": "Workbook has no sheets"})
    else:
        prs = Presentation(str(path))
        if len(prs.slides) == 0:
            issues.append({"code": "empty", "message": "Presentation has no slides"})
    return issues


def view_text(path: str | Path) -> str:
    """Alias of :func:`officekit.core.extract_text` for the ``view text`` CLI."""
    from .core import extract_text

    return extract_text(path)
