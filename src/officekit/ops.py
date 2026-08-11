"""Path-based DOM operations — the read/write layer agents can target.

Paths use 1-based indexing, mirroring agent-facing Office CLIs::

    docx : /paragraphs[1]   /headings[2]   /tables[1]/rows[1]/cells[2]
    xlsx : /sheets[1]        /sheets[1]/A1  /sheets[Name]!/A1
    pptx : /slides[1]/shapes[2]

A *step* is ``name[index]``; the xlsx final step may be a bare cell reference
(``A1``). ``get`` returns a JSON-friendly description; ``set_value`` writes to
the addressed element; ``add``/``remove`` mutate the document in place.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any, List, Tuple

from docx import Document as DocxDocument
from openpyxl import Workbook, load_workbook
from pptx import Presentation

from .core import _detect_kind

_PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE value, kept numeric to avoid extra deps


def _parse_path(path: str) -> List[Any]:
    steps: List[Any] = []
    for part in path.strip("/").split("/"):
        part = part.strip()
        if not part:
            continue
        m = re.match(r"^([A-Za-z_]+)\[(\d+)\]$", part)
        if m:
            steps.append((m.group(1), int(m.group(2)) - 1))  # 0-based
        else:
            # bare cell reference for xlsx, e.g. A1
            steps.append(("cell", part.upper()))
    if not steps:
        raise ValueError(f"Invalid path: {path!r}")
    return steps


def _navigate(root: Any, steps: List[Any]) -> Any:
    obj = root
    for name, sel in steps:
        if name == "cell":
            obj = obj[sel]
        elif name in ("sheets",):
            obj = obj.worksheets[sel]
        elif name in ("slides",):
            obj = obj.slides[sel]
        elif name in ("paragraphs", "headings"):
            obj = obj.paragraphs[sel]
        elif name in ("shapes",):
            obj = obj.shapes[sel]
        elif name in ("tables",):
            obj = obj.tables[sel]
        elif name in ("rows",):
            obj = obj.rows[sel]
        elif name in ("cells",):
            obj = obj.cells[sel]
        else:
            raise KeyError(f"Unknown path segment: {name!r}")
    return obj


def _load_root(path: str | Path) -> Tuple[Any, str]:
    kind = _detect_kind(path)
    if kind == "docx":
        return DocxDocument(str(path)), kind
    if kind == "xlsx":
        return load_workbook(str(path)), kind
    return Presentation(str(path)), kind


def _describe(obj: Any) -> Any:
    # xlsx cell
    if hasattr(obj, "value") and hasattr(obj, "coordinate"):
        return obj.value
    if hasattr(obj, "text_frame"):
        return obj.text_frame.text
    if hasattr(obj, "text"):
        return obj.text
    if hasattr(obj, "rows") and hasattr(obj, "cells"):
        return [[c.text for c in row.cells] for row in obj.rows]
    return str(obj)


# --------------------------------------------------------------------------
# get
# --------------------------------------------------------------------------
def get(path: str | Path, target: str) -> Any:
    """Resolve *target* (a path) inside the document at *path* and describe it."""
    root, _ = _load_root(path)
    obj = _navigate(root, _parse_path(target))
    return _describe(obj)


# --------------------------------------------------------------------------
# set_value
# --------------------------------------------------------------------------
def set_value(path: str | Path, target: str, value: str) -> None:
    """Write *value* to the element addressed by *target*, then save."""
    root, kind = _load_root(path)
    obj = _navigate(root, _parse_path(target))

    if kind == "xlsx" and hasattr(obj, "value") and hasattr(obj, "coordinate"):
        obj.value = value
    elif kind == "pptx" and hasattr(obj, "text_frame"):
        obj.text_frame.text = value
    elif kind == "docx":
        if hasattr(obj, "text") and hasattr(obj, "runs"):  # paragraph
            if obj.runs:
                obj.runs[0].text = value
                for r in obj.runs[1:]:
                    r.text = ""
            else:
                obj.add_run(value)
        elif hasattr(obj, "text"):  # table cell
            obj.text = value
        else:
            raise TypeError("Target is not a settable text element")
    else:
        raise TypeError("Target is not a settable text element")

    root.save(str(path))


# --------------------------------------------------------------------------
# add
# --------------------------------------------------------------------------
def add(path: str | Path, what: str, value: str | None = None) -> None:
    """Append a new element. *what* is one of: paragraph, slide, row, sheet."""
    root, kind = _load_root(path)
    if kind == "docx":
        if what == "paragraph":
            root.add_paragraph(value or "")
        elif what == "slide":
            raise ValueError("docx has no slides; use 'paragraph'")
        else:
            raise ValueError(f"docx add supports: paragraph")
    elif kind == "xlsx":
        if what == "row":
            ws = root.active
            ws.append(value.split("\t") if value else [])
        elif what == "sheet":
            root.create_sheet(title=value or None)
        else:
            raise ValueError("xlsx add supports: row, sheet")
    else:  # pptx
        if what == "slide":
            slide = root.slides.add_slide(root.slide_layouts[1])
            if slide.shapes.title is not None and value:
                slide.shapes.title.text = value
        else:
            raise ValueError("pptx add supports: slide")
    root.save(str(path))


# --------------------------------------------------------------------------
# remove
# --------------------------------------------------------------------------
def remove(path: str | Path, target: str) -> None:
    """Delete the element addressed by *target* (docx paragraph / pptx shape)."""
    root, kind = _load_root(path)
    obj = _navigate(root, _parse_path(target))
    el = getattr(obj, "_element", None)
    parent = el.getparent() if el is not None else None
    if parent is None:
        raise ValueError("Cannot remove this element")
    parent.remove(el)
    root.save(str(path))
